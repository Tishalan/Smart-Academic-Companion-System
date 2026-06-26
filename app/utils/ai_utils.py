import os
import speech_recognition as sr
from moviepy import VideoFileClip
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_JUSTIFY
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from datetime import datetime
import google.generativeai as genai
from pptx import Presentation
import fitz  # PyMuPDF
import json
import tempfile
import time
from sklearn.feature_extraction.text import TfidfVectorizer
import networkx as nx
import numpy as np
import re
import requests

class AIProcessor:
    def __init__(self):
        self.recognizer = sr.Recognizer()
        # Initialize Gemini if API key is present
        self.api_key = os.environ.get('GEMINI_API_KEY')
        self.openrouter_key = os.environ.get('OPENROUTER_API_KEY')
        self.provider = os.environ.get('AI_PROVIDER', 'gemini')
        
        if self.api_key:
            genai.configure(api_key=self.api_key)
            self.model = genai.GenerativeModel('gemini-pro')
        else:
            self.model = None

    def call_ai_api(self, prompt, deep_mode=False):
        """Unified method to call available AI providers."""
        # 1. Try OpenRouter (High quality fallback)
        if self.openrouter_key:
            try:
                response = requests.post(
                    url="https://openrouter.ai/api/v1/chat/completions",
                    headers={"Authorization": f"Bearer {self.openrouter_key}"},
                    json={
                        "model": "google/gemini-2.0-flash-exp:free",
                        "messages": [{"role": "user", "content": prompt}]
                    },
                    timeout=60
                )
                data = response.json()
                return data['choices'][0]['message']['content']
            except: pass
            
        # 2. Try Gemini (Primary)
        if self.model:
            try:
                response = self.model.generate_content(prompt)
                return response.text
            except: pass
            
        return None

    def speech_to_text(self, audio_path):
        """Convert speech in audio file to text using OpenAI Whisper (with Google fallback)."""
        # 1. Try Whisper (Professional Local Model as described in Radford et al., 2022)
        try:
            import whisper
            # Use 'base' model for optimal balance between speed and academic precision
            model = whisper.load_model("base")
            result = model.transcribe(audio_path)
            return result["text"]
        except Exception as e:
            # 2. Fallback to Google Speech Recognition if Whisper/FFmpeg fails
            print(f"Whisper Error (Concept Mode), falling back to Google: {str(e)}")
            try:
                with sr.AudioFile(audio_path) as source:
                    audio_data = self.recognizer.record(source)
                    text = self.recognizer.recognize_google(audio_data)
                    return text
            except Exception as ge:
                return f"Transcription Error: {str(ge)}"

    def get_ai_summary(self, text):
        """Generate a summary using Gemini."""
        if not self.model:
            return "AI Summary unavailable (API Key missing). Original Transcript: " + text[:500] + "..."
        
        try:
            prompt = f"Summarize the following lecture transcript into key learning points and a concise overview:\n\n{text}"
            response = self.model.generate_content(prompt)
            return response.text
        except Exception as e:
            return f"Error generating summary: {str(e)}"

    def extract_text_from_ppt(self, ppt_path):
        """Extract text from each slide of a PPT file."""
        try:
            prs = Presentation(ppt_path)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text)
                slides_text.append({
                    "slide_number": i + 1,
                    "content": "\n".join(slide_content)
                })
            return slides_text
        except Exception as e:
            raise Exception(f"PPT Error: {str(e)}")

    def extract_text_from_pdf(self, pdf_path):
        """Extract text from each page of a PDF file."""
        try:
            doc = fitz.open(pdf_path)
            pages_text = []
            for i, page in enumerate(doc):
                pages_text.append({
                    "slide_number": i + 1,
                    "content": page.get_text()
                })
            return pages_text
        except Exception as e:
            raise Exception(f"PDF Error: {str(e)}")

    def generate_detailed_study_notes(self, slides_data, module_name, deep_mode=True):
        """Use AI to deep-dive into each slide's concept with Wiki-Deep fallback."""
        detailed_notes = []
        
        for slide in slides_data:
            content = slide['content'].strip()
            if not content: continue
            
            # Use unified academic prompt
            prompt = (
                f"You are an expert academic tutor for '{module_name}'. "
                f"Explain the following lecture slide content in extreme detail. "
                f"Define all technical terms, provide examples, and ensure the student perfectly understands the concept.\n\n"
                f"SLIDE {slide['slide_number']} CONTENT:\n{content}"
            )
            
            result = self.call_ai_api(prompt, deep_mode)
            
            if not result:
                # Local Wiki-Deep Fallback if AI fails
                result = self.local_summarize(content)
                
            detailed_notes.append({
                "slide_number": slide['slide_number'],
                "original": content,
                "expansion": result
            })
            
        return detailed_notes

    def generate_pdf_report(self, notes_data, output_path, title, subtitle):
        """Generate a professional PDF with detailed notes as per user precision."""
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Heading 14, Bold, Times New Roman
        heading_style = ParagraphStyle(
            'CustomHeading',
            fontName='Times-Bold',
            fontSize=14,
            spaceBefore=15,
            spaceAfter=10,
            alignment=0, # Left aligned for heading or centered? Keeping it standard.
            textColor=colors.black
        )
        
        # Body 12, Justified, Line Spacing 1.5, Times New Roman
        body_style = ParagraphStyle(
            'CustomBody',
            fontName='Times-Roman',
            fontSize=12,
            leading=18, # 12 * 1.5
            alignment=TA_JUSTIFY,
            spaceAfter=12,
            textColor=colors.black
        )

        elements = []
        
        # Title page bit
        elements.append(Paragraph(title, heading_style))
        elements.append(Paragraph(subtitle, body_style))
        elements.append(Spacer(1, 20))
        
        if isinstance(notes_data, str):
            # If it's a single raw string (e.g. from Video deep mode)
            for line in notes_data.split('\n'):
                line = line.strip()
                if not line: continue
                # Basic markdown header detection
                if line.startswith('###'):
                    elements.append(Paragraph(line.replace('###', '').strip(), heading_style))
                elif line.startswith('##'):
                    elements.append(Paragraph(line.replace('##', '').strip(), heading_style))
                elif line.startswith('#'):
                    elements.append(Paragraph(line.replace('#', '').strip(), heading_style))
                else:
                    elements.append(Paragraph(line, body_style))
        else:
            # List of slide/section objects
            for note in notes_data:
                elements.append(Paragraph(f"Section {note.get('slide_number', '')}: Concept Analysis", heading_style))
                # Process the expansion text to handle basic formatting
                expansion = note['expansion']
                for line in expansion.split('\n'):
                    line = line.strip()
                    if not line: continue
                    if line.startswith('###') or line.isupper() or (line.startswith('**') and line.endswith('**')):
                        elements.append(Paragraph(line.replace('*', '').replace('#', '').strip(), heading_style))
                    else:
                        elements.append(Paragraph(line, body_style))
                elements.append(Spacer(1, 15))
            
        doc.build(elements)
        return output_path

    def research_concept(self, term):
        """Fetch detailed academic summary from Wikipedia for a term."""
        try:
            # Format term for URL
            formatted_term = term.strip().replace(' ', '_')
            url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{formatted_term}"
            response = requests.get(url, headers={'User-Agent': 'SmartAcademicSystem/1.0 (academic research tool)'}, timeout=5)
            
            if response.status_code == 200:
                data = response.json()
                extract = data.get('extract', '')
                if extract:
                    return {
                        'title': data.get('title', term),
                        'content': extract,
                        'url': data.get('content_urls', {}).get('desktop', {}).get('page', '')
                    }
            return None
        except:
            return None

    def extract_keywords(self, text):
        """Identify technical terms in a transcript (Simple Noun Phrase heuristic)."""
        # Common technical words to always look for if present
        tech_vocab = {
            'artificial intelligence', 'ai', 'generative ai', 'machine learning', 'deep learning',
            'neural network', 'transformer', 'algorithm', 'database', 'software engineering',
            'programming', 'data science', 'cloud computing', 'cybersecurity'
        }
        
        found = []
        lower_text = text.lower()
        for term in tech_vocab:
            if term in lower_text:
                found.append(term)
        
        # Also look for capitalized words that might be unique terms (excluding start of sentences)
        words = text.split()
        for i in range(1, len(words)):
            if words[i][0].isupper() and words[i-1][-1] not in '.!?':
                clean_word = re.sub(r'[^a-zA-Z]', '', words[i])
                if len(clean_word) > 3 and clean_word.lower() not in found:
                    found.append(clean_word)
                    
        return list(set(found))[:3] # Limit to top 3 for speed

    def local_summarize(self, text):
        """Enhanced Local Summarizer with Research-Mode (Wiki-Deep)."""
        # 1. Split into sentences
        sentences = re.split(r'(?<=[.!?]) +', text.strip())
        
        # 2. Extract Keywords for Deep Dive
        keywords = self.extract_keywords(text)
        
        # 3. Generate Research Content
        research_sections = []
        for kw in keywords:
            details = self.research_concept(kw)
            if details:
                research_sections.append(details)
        
        try:
            # Base logic same as before (TextRank)
            vectorizer = TfidfVectorizer(stop_words='english')
            tfidf_matrix = vectorizer.fit_transform(sentences)
            similarity_matrix = (tfidf_matrix * tfidf_matrix.T).toarray()
            nx_graph = nx.from_numpy_array(similarity_matrix)
            scores = nx.pagerank(nx_graph)
            ranked_sentences = sorted(((scores[i], s) for i, s in enumerate(sentences)), reverse=True)
            top_sentences = sorted([s for score, s in ranked_sentences[:min(len(sentences), 5)]], key=lambda x: sentences.index(x))
            
            # Format as Deep Academic Notes
            notes = "### LECTURE OVERVIEW (Local Summary)\n\n"
            for s in top_sentences:
                notes += f"* {s.strip()}\n"
            
            if research_sections:
                notes += "\n\n--- \n## DEEP DIVE RESEARCH (Academic Expansion)\n"
                for res in research_sections:
                    notes += f"\n### Concept: {res['title']}\n"
                    notes += f"{res['content']}\n"
                    if res['url']:
                        notes += f"\n*Referenced from: {res['url']}*\n"
            
            # Add heuristic Quiz
            notes += "\n\n--- \n## SELF-TEST QUIZ (Heuristic Mode)\n"
            if research_sections:
                notes += f"1. Based on the concept of '{research_sections[0]['title']}', how is it defined in modern academia?\n"
                notes += f"2. Compare the overview provided in the lecture with the technical details of {research_sections[0]['title']}.\n"
            else:
                notes += "1. What is the primary technical concept mentioned in this lecture?\n"
                notes += "2. How would you apply the discussed principles to a real-world project?\n"
            
            notes += "\n\n*Note: This deep guide was generated using Research-Mode (Free Academic API Fallback).* "
            return notes
        except Exception as e:
            return f"### RAW TRANSCRIPT\n\n{text}\n\n(Error in research mode: {str(e)})"

    def video_to_transcript(self, video_path, output_dir, deep_mode=False):
        """Full pipeline: Video -> Audio -> Text -> AI Notes -> PDF."""
        temp_audio = None
        try:
            # 1. Extract Audio to Temporary WAV
            video_filename = os.path.basename(video_path).split('.')[0]
            temp_audio = os.path.join(tempfile.gettempdir(), f"{video_filename}_{int(time.time())}.wav")
            
            video = VideoFileClip(video_path)
            # Use faster preset and optimized audio settings - moviepy 2.x uses logger instead of verbose
            video.audio.write_audiofile(temp_audio, codec='pcm_s16le', fps=16000, logger=None)
            video.close()
            
            # 2. Transcribe Audio
            raw_text = self.speech_to_text(temp_audio)
            
            if "Error processing speech" in raw_text:
                return raw_text
                
            # 3. AI Smart Notes Generation
            prompt_instr = "provide extreme detail on every technical concept mentioned. Even if just mentioned in passing, explain its history, importance, and mechanics. ADD A 'SELF-TEST QUIZ' at the end." if deep_mode else "Convert the following raw lecture transcript into detailed, professional study notes."
            
            prompt = (
                f"You are a Senior Academic Researcher and expert scribe. \n\n"
                f"TRANSCRIPT:\n{raw_text}\n\n"
                f"INSTRUCTIONS:\n{prompt_instr}"
            )
            
            formatted_notes = self.call_ai_api(prompt, deep_mode)
            
            if not formatted_notes:
                # Local research mode fallback
                formatted_notes = self.local_summarize(raw_text)
                if deep_mode:
                    formatted_notes += "\n\n(Note: Research-Mode Fallback used as AI APIs were unavailable)."
            
            # 4. Generate PDF with specific user requirements
            pdf_filename = f"Lecture_Notes_{video_filename}_{int(time.time())}.pdf"
            pdf_path = os.path.join(output_dir, pdf_filename)
            
            title_prefix = "DEEP Academic Study Guide" if deep_mode else "Lecture Notes"
            self.generate_pdf_report(
                notes_data=formatted_notes,
                output_path=pdf_path,
                title=f"{title_prefix}: {video_filename}",
                subtitle=f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            )
            
            return pdf_path
            
        except Exception as e:
            return f"Processing Error: {str(e)}"
        finally:
            if temp_audio and os.path.exists(temp_audio):
                try: os.remove(temp_audio)
                except: pass

ai_processor = AIProcessor()
